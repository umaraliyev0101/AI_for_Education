"""
Presentation Processing Service
================================

Extracts text from presentations (PPTX/PDF), generates TTS audio
for each slide, and prepares structured data for frontend display.
"""

import os
import logging
from typing import List, Dict, Optional, Any
from pathlib import Path
import json
import asyncio
from PIL import Image
import io
import functools

logger = logging.getLogger(__name__)

# Import presentation processing libraries
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("PyPDF2 not available")

# For PPTX to image conversion on Windows
try:
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    COMTYPES_AVAILABLE = False
    logger.warning("⚠️ comtypes not available - PPTX image conversion may be limited")

# For PDF to image conversion
try:
    from pdf2image import convert_from_path
    try:
        # specific exception used when poppler's pdfinfo is missing
        from pdf2image.exceptions import PDFInfoNotInstalledError
    except Exception:
        PDFInfoNotInstalledError = None
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("⚠️ pdf2image not available - PDF image conversion disabled")

# For PDF to image conversion (fallback, no external dependencies)
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.warning("⚠️ PyMuPDF not available - PDF fallback conversion disabled")

try:
    from stt_pipelines.uzbek_tts_pipeline import create_uzbek_tts
    TTS_AVAILABLE = True
except ImportError:
    TTS_AVAILABLE = False
    logger.warning("TTS pipeline not available")

# Check for LibreOffice availability (fallback for PPTX conversion)
LIBREOFFICE_AVAILABLE = False
try:
    import subprocess
    result = subprocess.run(['soffice', '--version'], capture_output=True, text=True, timeout=5)
    if result.returncode == 0:
        LIBREOFFICE_AVAILABLE = True
        logger.info("✅ LibreOffice available for PPTX conversion fallback")
    else:
        logger.warning("⚠️ LibreOffice soffice command not found or not working")
except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
    logger.warning(f"⚠️ LibreOffice check failed: {e}")
    logger.warning("💡 Install LibreOffice for PPTX conversion fallback on systems without PowerPoint")


class PresentationService:
    """Service for processing presentations and generating audio"""
    
    def __init__(
        self, 
        audio_output_dir: str = "./uploads/audio/presentations",
        slides_output_dir: str = "./uploads/slides"
    ):
        self.audio_output_dir = audio_output_dir
        self.slides_output_dir = slides_output_dir
        os.makedirs(audio_output_dir, exist_ok=True)
        os.makedirs(slides_output_dir, exist_ok=True)
        
        # Initialize TTS
        self.tts = None
        if TTS_AVAILABLE:
            try:
                self.tts = create_uzbek_tts(voice="male_clear")
                logger.info("✅ TTS initialized for presentations")
            except Exception as e:
                logger.error(f"❌ Failed to initialize TTS: {e}")
    
    def _convert_pptx_to_images(self, pptx_path: str, lesson_id: int) -> List[str]:
        """
        Convert PPTX slides to PNG images using PowerPoint COM or LibreOffice fallback
        
        Returns:
            List of image file paths
        """
        # Create output directory for this lesson
        lesson_slides_dir = os.path.join(self.slides_output_dir, f"lesson_{lesson_id}")
        os.makedirs(lesson_slides_dir, exist_ok=True)
        
        logger.info(f"🖼️  Slide directory: {os.path.abspath(lesson_slides_dir)}")
        
        # ✅ Check if images already exist
        existing_images = []
        try:
            for file in os.listdir(lesson_slides_dir):
                if file.startswith('slide_') and file.endswith('.png'):
                    rel_path = f"uploads/slides/lesson_{lesson_id}/{file}"
                    existing_images.append((int(file.split('_')[1].split('.')[0]), rel_path))
        except Exception as e:
            logger.warning(f"⚠️  Error checking existing images: {e}")
        
        if existing_images:
            # Sort by slide number and return paths
            existing_images.sort(key=lambda x: x[0])
            image_paths = [path for _, path in existing_images]
            logger.info(f"✅ Found {len(image_paths)} existing slide images")
            return image_paths
        
        # Try PowerPoint COM first (Windows with PowerPoint)
        if COMTYPES_AVAILABLE:
            logger.info("🔄 Attempting PowerPoint COM conversion...")
            powerpoint_result = self._convert_pptx_with_powerpoint(pptx_path, lesson_id, lesson_slides_dir)
            if powerpoint_result:
                return powerpoint_result
        
        # Fallback to LibreOffice (cross-platform)
        # Probe for LibreOffice at runtime and attempt conversion even if the
        # import-time check failed (some environments change PATH after startup).
        logger.info("🔄 Attempting LibreOffice conversion fallback (runtime probe)...")
        if not PDF2IMAGE_AVAILABLE:
            logger.info("ℹ️ pdf2image not available; will attempt LibreOffice -> PyMuPDF path")
        libreoffice_result = self._convert_pptx_with_libreoffice(pptx_path, lesson_id, lesson_slides_dir)
        if libreoffice_result:
            return libreoffice_result
        
        # Both methods failed
        logger.error("❌ All PPTX conversion methods failed")
        # Give actionable hints
        if not COMTYPES_AVAILABLE:
            logger.error("💡 PowerPoint COM not available: pip install comtypes and ensure MS PowerPoint is installed")
        if not LIBREOFFICE_AVAILABLE:
            logger.error("💡 LibreOffice not found: install LibreOffice or ensure 'soffice' is in PATH")
        if not PDF2IMAGE_AVAILABLE:
            logger.error("💡 pdf2image not installed: pip install pdf2image (requires poppler) or rely on PyMuPDF fallback")
        if not PYMUPDF_AVAILABLE:
            logger.error("💡 PyMuPDF not installed: pip install PyMuPDF")
        return []
    
    def _convert_pptx_with_powerpoint(self, pptx_path: str, lesson_id: int, output_dir: str) -> List[str]:
        """Convert PPTX to images using PowerPoint COM"""
        try:
            logger.info(f"🚀 Starting PowerPoint COM automation for {pptx_path}")
            
            # Convert to absolute path
            abs_pptx_path = os.path.abspath(pptx_path)
            abs_output_dir = os.path.abspath(output_dir)
            
            logger.info(f"📂 Input: {abs_pptx_path}")
            logger.info(f"📂 Output: {abs_output_dir}")
            
            # Verify input file exists
            if not os.path.exists(abs_pptx_path):
                logger.error(f"❌ Input file not found: {abs_pptx_path}")
                return []
            
            # Initialize PowerPoint COM
            logger.info("🔌 Creating PowerPoint COM object...")
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            
            # ✅ Configure PowerPoint for automation
            # Note: PowerPoint 2016+ does NOT allow hiding via Visible=0
            # The window will appear briefly during export - this is normal
            try:
                powerpoint.DisplayAlerts = 0  # Disable all alerts/dialogs (ppAlertsNone)
            except Exception as e:
                logger.warning(f"⚠️  Could not set DisplayAlerts: {e}")
            
            logger.info("📖 Opening presentation...")
            # Open presentation without displaying window (WithWindow=False)
            presentation = powerpoint.Presentations.Open(
                abs_pptx_path,
                ReadOnly=True,
                Untitled=False,
                WithWindow=False
            )
            
            slide_count = presentation.Slides.Count
            logger.info(f"📊 Found {slide_count} slides to export")
            
            image_paths = []
            
            # Export each slide as PNG
            for i in range(1, slide_count + 1):
                image_filename = f"slide_{i}.png"
                image_path = os.path.join(abs_output_dir, image_filename)
                
                logger.info(f"🖼️  Exporting slide {i}/{slide_count}...")
                
                # Export slide (2 = ppSaveAsPNG)
                presentation.Slides(i).Export(image_path, "PNG")
                
                # Verify file was created
                if os.path.exists(image_path):
                    # Store relative path
                    rel_path = f"uploads/slides/lesson_{lesson_id}/{image_filename}"
                    image_paths.append(rel_path)
                    logger.info(f"✅ Exported slide {i} to {image_filename} ({os.path.getsize(image_path)} bytes)")
                else:
                    logger.error(f"❌ Failed to export slide {i} - file not created")
            
            # Close presentation
            logger.info("🔒 Closing presentation...")
            presentation.Close()
            
            logger.info("🔌 Quitting PowerPoint...")
            powerpoint.Quit()
            
            logger.info(f"✅ Successfully exported {len(image_paths)} slides with PowerPoint")
            return image_paths
            
        except Exception as e:
            logger.error(f"❌ PowerPoint COM conversion failed: {e}")
            logger.error(f"❌ Error type: {type(e).__name__}")
            logger.error(f"❌ Error details: {str(e)}")
            import traceback
            logger.error(f"❌ Traceback:\n{traceback.format_exc()}")
            return []
    
    def _convert_pptx_with_libreoffice(self, pptx_path: str, lesson_id: int, output_dir: str) -> List[str]:
        """Convert PPTX to images using LibreOffice (fallback method)"""
        try:
            logger.info(f"📄 Converting PPTX with LibreOffice: {pptx_path}")
            
            # Convert to absolute paths
            abs_pptx_path = os.path.abspath(pptx_path)
            abs_output_dir = os.path.abspath(output_dir)
            
            # Create temporary PDF path
            temp_pdf_path = os.path.join(abs_output_dir, f"temp_presentation.pdf")
            
            logger.info(f"📂 Input: {abs_pptx_path}")
            logger.info(f"📂 Temp PDF: {temp_pdf_path}")
            logger.info(f"📂 Output dir: {abs_output_dir}")
            
            # DEBUG: Log environment information
            logger.info(f"🐛 DEBUG: Current PATH: {os.environ.get('PATH', 'NOT_SET')}")
            logger.info(f"🐛 DEBUG: Working directory: {os.getcwd()}")
            
            # Try to find soffice executable
            import shutil
            # Allow explicit env override for the soffice binary (helpful for services)
            soffice_env = os.environ.get('LIBREOFFICE_PATH') or os.environ.get('SOFFICE_PATH')
            if soffice_env:
                # If user provides a folder, append executable name
                candidate = os.path.expandvars(soffice_env)
                if os.path.isdir(candidate):
                    candidate = os.path.join(candidate, 'soffice.exe' if os.name == 'nt' else 'soffice')
                soffice_path = candidate if os.path.exists(candidate) else None
                logger.info(f"🐛 DEBUG: Using soffice from env var: {soffice_env} -> {soffice_path}")
            else:
                soffice_path = shutil.which("soffice")
                logger.info(f"🐛 DEBUG: shutil.which('soffice'): {soffice_path}")

            if not soffice_path:
                # Try common LibreOffice installation locations
                common_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    r"C:\Users\%USERNAME%\AppData\Local\Programs\LibreOffice\program\soffice.exe",
                    "/usr/bin/soffice",
                    "/usr/lib/libreoffice/program/soffice",
                    "/opt/libreoffice/program/soffice"
                ]
                
                for path in common_paths:
                    expanded_path = os.path.expandvars(path)
                    logger.info(f"🐛 DEBUG: Checking path: {expanded_path}")
                    if os.path.exists(expanded_path):
                        soffice_path = expanded_path
                        logger.info(f"🐛 DEBUG: Found soffice at: {soffice_path}")
                        break
            
            if not soffice_path:
                logger.error("❌ soffice executable not found in PATH or common locations")
                logger.error("💡 Please ensure LibreOffice is installed and 'soffice' is in PATH")
                logger.error("💡 Or set an explicit environment variable before starting the backend:")
                logger.error("  - PowerShell (current session): $env:LIBREOFFICE_PATH='C:\\Program Files\\LibreOffice\\program'")
                logger.error("  - Windows (permanent): setx LIBREOFFICE_PATH 'C:\\Program Files\\LibreOffice\\program'")
                logger.error("  - Linux: export LIBREOFFICE_PATH=/usr/bin")
                return []
            
            # Verify input file exists
            if not os.path.exists(abs_pptx_path):
                logger.error(f"❌ Input file not found: {abs_pptx_path}")
                return []
            
            # Step 1: Convert PPTX to PDF using LibreOffice
            logger.info("🔄 Running LibreOffice conversion...")
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', abs_output_dir,
                abs_pptx_path
            ]
            
            logger.info(f"🐛 DEBUG: Command: {' '.join(cmd)}")
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            logger.info(f"🐛 DEBUG: Return code: {result.returncode}")
            logger.info(f"🐛 DEBUG: stdout: {result.stdout}")
            if result.stderr:
                logger.warning(f"🐛 DEBUG: stderr: {result.stderr}")

            if result.returncode != 0:
                logger.error(f"❌ LibreOffice conversion failed: {result.stderr}")
                return []

            # Detect generated PDF file(s).
            try:
                files_after = [f for f in os.listdir(abs_output_dir) if f.lower().endswith('.pdf')]
            except Exception:
                files_after = []

            # Prefer file with same stem as source PPTX
            base_stem = os.path.splitext(os.path.basename(pptx_path))[0]
            matching = [f for f in files_after if os.path.splitext(f)[0] == base_stem]

            chosen_pdf = None
            if matching:
                chosen_pdf = os.path.join(abs_output_dir, matching[0])
                logger.info(f"🐛 DEBUG: Found PDF matching PPTX stem: {matching[0]}")
            elif files_after:
                # If there are multiple PDFs, pick the newest by mtime
                files_full = [os.path.join(abs_output_dir, f) for f in files_after]
                files_full.sort(key=lambda p: os.path.getmtime(p), reverse=True)
                chosen_pdf = files_full[0]
                logger.info(f"🐛 DEBUG: No stem match; picking newest PDF: {os.path.basename(chosen_pdf)}")

            if not chosen_pdf:
                # Fallback: maybe LibreOffice produced a file with a different name
                logger.error("❌ PDF file not found after LibreOffice conversion")
                logger.error(f"🐛 DEBUG: Files in output dir: {os.listdir(abs_output_dir)}")
                return []

            temp_pdf_path = chosen_pdf
            logger.info(f"✅ PDF created: {temp_pdf_path}")
            
            # Step 2: Convert PDF to images using pdf2image
            logger.info("🖼️ Converting PDF to images...")
            # Allow overriding poppler bin path via env var POPPLER_PATH
            poppler_path = os.environ.get('POPPLER_PATH')
            try:
                if poppler_path:
                    logger.info(f"🐛 DEBUG: Using POPPLER_PATH={poppler_path}")
                    images = convert_from_path(temp_pdf_path, dpi=150, poppler_path=poppler_path)
                else:
                    images = convert_from_path(temp_pdf_path, dpi=150)
            except Exception as e:
                # If pdf2image fails because poppler/pdfinfo is missing, fall back to PyMuPDF
                err_name = type(e).__name__
                logger.warning(f"⚠️ pdf2image failed: {err_name}: {e}")
                try:
                    from pdf2image.exceptions import PDFInfoNotInstalledError
                except Exception:
                    PDFInfoNotInstalledError = None

                if PDFInfoNotInstalledError and isinstance(e, PDFInfoNotInstalledError):
                    logger.info("🔁 Falling back to PyMuPDF for PDF->PNG conversion")
                    return self._convert_pdf_with_pymupdf(temp_pdf_path, lesson_id, abs_output_dir)
                else:
                    logger.error(f"❌ pdf2image conversion failed: {e}")
                    return []
            
            image_paths = []
            
            for i, image in enumerate(images, start=1):
                image_filename = f"slide_{i}.png"
                image_path = os.path.join(abs_output_dir, image_filename)
                
                # Save image
                image.save(image_path, 'PNG')
                
                # Store relative path
                rel_path = f"uploads/slides/lesson_{lesson_id}/{image_filename}"
                image_paths.append(rel_path)
                
                logger.info(f"✅ Exported slide {i} to {image_filename}")
            
            # Clean up temporary PDF
            try:
                if os.path.exists(temp_pdf_path):
                    os.remove(temp_pdf_path)
                    logger.info("🧹 Cleaned up temporary PDF file")
            except Exception as e:
                logger.warning(f"⚠️ Failed to clean up temp PDF: {e}")
            
            logger.info(f"✅ Successfully exported {len(image_paths)} slides with LibreOffice")
            return image_paths
            
        except subprocess.TimeoutExpired:
            logger.error("❌ LibreOffice conversion timed out")
            return []
        except Exception as e:
            logger.error(f"❌ LibreOffice conversion failed: {e}")
            logger.error(f"❌ Error type: {type(e).__name__}")
            import traceback
            logger.error(f"❌ Traceback:\n{traceback.format_exc()}")
            return []
    
    def _convert_pdf_to_images(self, pdf_path: str, lesson_id: int) -> List[str]:
        """
        Convert PDF pages to PNG images using pdf2image or PyMuPDF fallback
        
        Returns:
            List of image file paths
        """
        # Create output directory for this lesson
        lesson_slides_dir = os.path.join(self.slides_output_dir, f"lesson_{lesson_id}")
        os.makedirs(lesson_slides_dir, exist_ok=True)
        
        logger.info(f"🖼️  PDF directory: {os.path.abspath(lesson_slides_dir)}")
        
        # ✅ Check if images already exist
        existing_images = []
        try:
            for file in os.listdir(lesson_slides_dir):
                if file.startswith('slide_') and file.endswith('.png'):
                    rel_path = f"uploads/slides/lesson_{lesson_id}/{file}"
                    existing_images.append((int(file.split('_')[1].split('.')[0]), rel_path))
        except Exception as e:
            logger.warning(f"⚠️  Error checking existing images: {e}")
        
        if existing_images:
            # Sort by slide number and return paths
            existing_images.sort(key=lambda x: x[0])
            image_paths = [path for _, path in existing_images]
            logger.info(f"✅ Found {len(image_paths)} existing PDF page images")
            return image_paths
        
        # Try pdf2image first (requires poppler)
        if PDF2IMAGE_AVAILABLE:
            logger.info("🔄 Attempting pdf2image conversion...")
            pdf2image_result = self._convert_pdf_with_pdf2image(pdf_path, lesson_id, lesson_slides_dir)
            if pdf2image_result:
                return pdf2image_result
        
        # Fallback to PyMuPDF (no external dependencies)
        if PYMUPDF_AVAILABLE:
            logger.info("🔄 Attempting PyMuPDF conversion fallback...")
            pymupdf_result = self._convert_pdf_with_pymupdf(pdf_path, lesson_id, lesson_slides_dir)
            if pymupdf_result:
                return pymupdf_result
        
        # Both methods failed
        logger.error("❌ All PDF conversion methods failed")
        if not PDF2IMAGE_AVAILABLE:
            logger.error("💡 Install pdf2image: pip install pdf2image")
            logger.error("💡 Also install poppler: conda install -c conda-forge poppler (or system package manager)")
        if not PYMUPDF_AVAILABLE:
            logger.error("💡 Install PyMuPDF: pip install PyMuPDF")
        return []
    
    def _convert_pdf_with_pdf2image(self, pdf_path: str, lesson_id: int, output_dir: str) -> List[str]:
        """Convert PDF to images using pdf2image (requires poppler)"""
        try:
            logger.info(f"📄 Converting PDF with pdf2image: {pdf_path}")
            
            # Convert PDF to images
            images = convert_from_path(pdf_path, dpi=150)
            
            image_paths = []
            
            for i, image in enumerate(images, start=1):
                image_filename = f"slide_{i}.png"
                image_path = os.path.join(output_dir, image_filename)
                
                # Save image
                image.save(image_path, 'PNG')
                
                # Store relative path
                rel_path = f"uploads/slides/lesson_{lesson_id}/{image_filename}"
                image_paths.append(rel_path)
                
                logger.info(f"✅ Exported page {i} to {image_filename}")
            
            logger.info(f"✅ Successfully converted {len(image_paths)} pages with pdf2image")
            return image_paths
            
        except Exception as e:
            logger.error(f"❌ pdf2image conversion failed: {e}")
            return []
    
    def _convert_pdf_with_pymupdf(self, pdf_path: str, lesson_id: int, output_dir: str) -> List[str]:
        """Convert PDF to images using PyMuPDF (no external dependencies)"""
        try:
            logger.info(f"📄 Converting PDF with PyMuPDF: {pdf_path}")
            
            # Open PDF document
            doc = fitz.open(pdf_path)
            image_paths = []
            
            # Convert each page to image
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Render page to image (higher resolution for better quality)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scaling for better quality
                
                image_filename = f"slide_{page_num + 1}.png"
                image_path = os.path.join(output_dir, image_filename)
                
                # Save image
                pix.save(image_path)
                
                # Store relative path
                rel_path = f"uploads/slides/lesson_{lesson_id}/{image_filename}"
                image_paths.append(rel_path)
                
                logger.info(f"✅ Exported page {page_num + 1} to {image_filename}")
            
            doc.close()
            logger.info(f"✅ Successfully converted {len(image_paths)} pages with PyMuPDF")
            return image_paths
            
        except Exception as e:
            logger.error(f"❌ PyMuPDF conversion failed: {e}")
            return []
    
    async def process_presentation(self, presentation_path: str, lesson_id: int) -> Optional[Dict[str, Any]]:
        """
        Process a presentation file and generate audio for each slide
        
        Args:
            presentation_path: Path to presentation file (PPTX or PDF)
            lesson_id: Lesson database ID
            
        Returns:
            Dictionary with slides and audio paths, or None on error
        """
        try:
            file_ext = Path(presentation_path).suffix.lower()
            
            if file_ext == '.pptx' and PPTX_AVAILABLE:
                return await self._process_pptx(presentation_path, lesson_id)
            elif file_ext == '.pdf' and PDF_AVAILABLE:
                return await self._process_pdf(presentation_path, lesson_id)
            else:
                logger.error(f"❌ Unsupported file format: {file_ext}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Failed to process presentation: {e}")
            return None
    
    async def process_presentation_with_progress(
        self, 
        presentation_path: str, 
        lesson_id: int,
        progress_callback=None
    ) -> Optional[Dict[str, Any]]:
        """
        Process presentation with real-time progress updates
        
        Args:
            presentation_path: Path to presentation file
            lesson_id: Lesson database ID
            progress_callback: Async function(current, total, text) for progress updates
            
        Returns:
            Dictionary with slides and audio paths
        """
        try:
            file_ext = Path(presentation_path).suffix.lower()
            
            if file_ext == '.pptx' and PPTX_AVAILABLE:
                return await self._process_pptx_with_progress(
                    presentation_path, lesson_id, progress_callback
                )
            elif file_ext == '.pdf' and PDF_AVAILABLE:
                return await self._process_pdf_with_progress(
                    presentation_path, lesson_id, progress_callback
                )
            else:
                logger.error(f"❌ Unsupported file format: {file_ext}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Failed to process presentation: {e}")
            return None
    
    async def _process_pptx(self, pptx_path: str, lesson_id: int) -> Dict[str, Any]:
        """Extract text, generate audio, and create slide images from PPTX"""
        prs = Presentation(pptx_path)
        
        # Convert slides to images first
        logger.info(f"🖼️ Converting PPTX to images...")
        image_paths = self._convert_pptx_to_images(pptx_path, lesson_id)
        
        slides_data = []
        
        for idx, slide in enumerate(prs.slides):
            slide_number = idx + 1
            
            # Extract text from slide
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            slide_text = "\n".join(text_content)
            
            if not slide_text:
                slide_text = f"Slayd {slide_number}"
            
            # Generate audio for slide
            audio_path = await self._generate_slide_audio(
                slide_text, 
                lesson_id, 
                slide_number
            )
            
            # Get image path for this slide
            image_path = image_paths[idx] if idx < len(image_paths) else None
            
            slides_data.append({
                'slide_number': slide_number,
                'text': slide_text,
                'audio_path': audio_path,
                'image_path': image_path,  # ← NEW: Include slide image
                'duration_estimate': self._estimate_audio_duration(slide_text)
            })
            
            logger.info(f"✅ Processed slide {slide_number}/{len(prs.slides)}")
        
        # Save presentation metadata
        presentation_data = {
            'lesson_id': lesson_id,
            'total_slides': len(slides_data),
            'slides': slides_data,
            'processed_at': str(asyncio.get_event_loop().time())
        }
        
        # Save metadata to JSON inside lesson folder
        lesson_audio_dir = os.path.join(self.audio_output_dir, f"lesson_{lesson_id}")
        os.makedirs(lesson_audio_dir, exist_ok=True)
        
        metadata_path = os.path.join(lesson_audio_dir, "metadata.json")
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, ensure_ascii=False, indent=2)
        
        presentation_data['metadata_path'] = metadata_path
        
        logger.info(f"✅ Processed {len(slides_data)} slides for lesson {lesson_id}")
        return presentation_data
    
    async def _process_pptx_with_progress(
        self, 
        pptx_path: str, 
        lesson_id: int,
        progress_callback=None
    ) -> Dict[str, Any]:
        """Extract text, generate audio, create images with progress updates"""
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        
        # Convert slides to images first
        logger.info(f"🖼️ Converting PPTX to images...")
        image_paths = self._convert_pptx_to_images(pptx_path, lesson_id)
        
        slides_data = []
        
        for idx, slide in enumerate(prs.slides):
            slide_number = idx + 1
            
            # Extract text from slide
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            slide_text = "\n".join(text_content)
            
            if not slide_text:
                slide_text = f"Slayd {slide_number}"
            
            # ✅ NEW: Send progress update
            if progress_callback:
                await progress_callback(slide_number, total_slides, slide_text)
            
            # Generate audio for slide
            audio_path = await self._generate_slide_audio(
                slide_text, 
                lesson_id, 
                slide_number
            )
            
            # Get image path for this slide
            image_path = image_paths[idx] if idx < len(image_paths) else None
            
            slides_data.append({
                'slide_number': slide_number,
                'text': slide_text,
                'audio_path': audio_path,
                'image_path': image_path,
                'duration_estimate': self._estimate_audio_duration(slide_text)
            })
            
            logger.info(f"✅ Processed slide {slide_number}/{total_slides}")
        
        # Save presentation metadata
        presentation_data = {
            'lesson_id': lesson_id,
            'total_slides': len(slides_data),
            'slides': slides_data,
            'processed_at': str(asyncio.get_event_loop().time())
        }
        
        # Save metadata to JSON inside lesson folder
        lesson_audio_dir = os.path.join(self.audio_output_dir, f"lesson_{lesson_id}")
        os.makedirs(lesson_audio_dir, exist_ok=True)
        
        metadata_path = os.path.join(lesson_audio_dir, "metadata.json")
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, ensure_ascii=False, indent=2)
        
        presentation_data['metadata_path'] = metadata_path
        
        logger.info(f"✅ Processed {len(slides_data)} slides for lesson {lesson_id}")
        return presentation_data
    
    async def _process_pdf(self, pdf_path: str, lesson_id: int) -> Dict[str, Any]:
        """Extract text, generate audio, and create page images from PDF"""
        # Convert PDF pages to images first
        logger.info(f"🖼️ Converting PDF to images...")
        image_paths = self._convert_pdf_to_images(pdf_path, lesson_id)
        
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            slides_data = []
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text().strip()
                
                if not page_text:
                    page_text = f"Sahifa {page_num + 1}"
                
                # Generate audio for page
                audio_path = await self._generate_slide_audio(
                    page_text,
                    lesson_id,
                    page_num + 1
                )
                
                # Get image path for this page
                image_path = image_paths[page_num] if page_num < len(image_paths) else None
                
                slides_data.append({
                    'slide_number': page_num + 1,
                    'text': page_text,
                    'audio_path': audio_path,
                    'image_path': image_path,  # ← NEW: Include page image
                    'duration_estimate': self._estimate_audio_duration(page_text)
                })
                
                logger.info(f"✅ Processed page {page_num + 1}/{len(pdf_reader.pages)}")
        
        # Save presentation metadata
        presentation_data = {
            'lesson_id': lesson_id,
            'total_slides': len(slides_data),
            'slides': slides_data,
            'processed_at': str(asyncio.get_event_loop().time())
        }
        
        # Save metadata to JSON inside lesson folder
        lesson_audio_dir = os.path.join(self.audio_output_dir, f"lesson_{lesson_id}")
        os.makedirs(lesson_audio_dir, exist_ok=True)
        
        metadata_path = os.path.join(lesson_audio_dir, "metadata.json")
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, ensure_ascii=False, indent=2)
        
        presentation_data['metadata_path'] = metadata_path
        
        logger.info(f"✅ Processed {len(slides_data)} pages for lesson {lesson_id}")
        return presentation_data
    
    async def _process_pdf_with_progress(
        self, 
        pdf_path: str, 
        lesson_id: int,
        progress_callback=None
    ) -> Dict[str, Any]:
        """Extract text, generate audio, create images with progress updates"""
        # Convert PDF pages to images first
        logger.info(f"🖼️ Converting PDF to images...")
        image_paths = self._convert_pdf_to_images(pdf_path, lesson_id)
        
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            slides_data = []
            
            for page_num in range(total_pages):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text().strip()
                
                if not page_text:
                    page_text = f"Sahifa {page_num + 1}"
                
                # ✅ NEW: Send progress update
                if progress_callback:
                    await progress_callback(page_num + 1, total_pages, page_text)
                
                # Generate audio for page
                audio_path = await self._generate_slide_audio(
                    page_text,
                    lesson_id,
                    page_num + 1
                )
                
                # Get image path for this page
                image_path = image_paths[page_num] if page_num < len(image_paths) else None
                
                slides_data.append({
                    'slide_number': page_num + 1,
                    'text': page_text,
                    'audio_path': audio_path,
                    'image_path': image_path,
                    'duration_estimate': self._estimate_audio_duration(page_text)
                })
                
                logger.info(f"✅ Processed page {page_num + 1}/{total_pages}")
        
        # Save presentation metadata
        presentation_data = {
            'lesson_id': lesson_id,
            'total_slides': len(slides_data),
            'slides': slides_data,
            'processed_at': str(asyncio.get_event_loop().time())
        }
        
        # Save metadata to JSON inside lesson folder
        lesson_audio_dir = os.path.join(self.audio_output_dir, f"lesson_{lesson_id}")
        os.makedirs(lesson_audio_dir, exist_ok=True)
        
        metadata_path = os.path.join(lesson_audio_dir, "metadata.json")
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, ensure_ascii=False, indent=2)
        
        presentation_data['metadata_path'] = metadata_path
        
        logger.info(f"✅ Processed {len(slides_data)} pages for lesson {lesson_id}")
        return presentation_data
    
    async def _generate_slide_audio(self, text: str, lesson_id: int, slide_number: int) -> Optional[str]:
        """Generate TTS audio for a slide"""
        if not self.tts:
            logger.warning("⚠️ TTS not available, skipping audio generation")
            return None
        
        try:
            # Create lesson-specific audio directory
            lesson_audio_dir = os.path.join(self.audio_output_dir, f"lesson_{lesson_id}")
            os.makedirs(lesson_audio_dir, exist_ok=True)
            
            # Create audio file path with simpler naming
            audio_filename = f"slide_{slide_number}.mp3"
            audio_path = os.path.join(lesson_audio_dir, audio_filename)
            
            # Generate speech WITHOUT playing (run in executor to avoid blocking)
            # Use the running loop and run in executor to avoid 'no current event loop' errors
            loop = asyncio.get_running_loop()
            speak_fn = functools.partial(self.tts.speak_text, text, save_to_file=audio_path, play_audio=False)
            logger.debug(f"🗣️ Generating TTS for lesson {lesson_id} slide {slide_number} -> {audio_path}")
            await loop.run_in_executor(None, speak_fn)
            logger.debug(f"🗣️ TTS generation completed for lesson {lesson_id} slide {slide_number}")
            
            if os.path.exists(audio_path):
                logger.info(f"✅ Generated audio for slide {slide_number}")
                # Return relative path with leading slash for frontend
                return f"/uploads/audio/presentations/lesson_{lesson_id}/{audio_filename}"
            else:
                logger.error(f"❌ Audio file not created for slide {slide_number} at expected path: {audio_path}")
                return None
                
        except Exception as e:
            logger.error(f"❌ Failed to generate audio for slide {slide_number}: {e}")
            return None
    
    def _estimate_audio_duration(self, text: str) -> float:
        """Estimate audio duration in seconds based on text length"""
        # Average reading speed: ~150 words per minute
        words = len(text.split())
        duration = (words / 150) * 60
        return round(duration, 1)
    
    def load_presentation_metadata(self, lesson_id: int) -> Optional[Dict[str, Any]]:
        """Load previously processed presentation metadata"""
        # Try new location first (inside lesson folder)
        lesson_audio_dir = os.path.join(self.audio_output_dir, f"lesson_{lesson_id}")
        metadata_path = os.path.join(lesson_audio_dir, "metadata.json")
        
        # Fallback to old location for backward compatibility
        if not os.path.exists(metadata_path):
            metadata_path = os.path.join(
                self.audio_output_dir,
                f"lesson_{lesson_id}_presentation_metadata.json"
            )
        
        if os.path.exists(metadata_path):
            try:
                with open(metadata_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    
                    # ✅ FIX: Ensure all paths have leading slash for frontend
                    if 'slides' in data:
                        for slide in data['slides']:
                            # Fix image_path
                            if slide.get('image_path') and not slide['image_path'].startswith('/'):
                                slide['image_path'] = '/' + slide['image_path']
                            
                            # Fix audio_path
                            if slide.get('audio_path') and not slide['audio_path'].startswith('/'):
                                slide['audio_path'] = '/' + slide['audio_path']
                    
                    return data
            except Exception as e:
                logger.error(f"❌ Failed to load metadata: {e}")
        
        return None
    
    def get_slide_data(self, lesson_id: int, slide_number: int) -> Optional[Dict[str, Any]]:
        """Get data for a specific slide"""
        metadata = self.load_presentation_metadata(lesson_id)
        if metadata and 'slides' in metadata:
            for slide in metadata['slides']:
                if slide['slide_number'] == slide_number:
                    return slide
        return None
    
    async def regenerate_slide_audio(self, lesson_id: int, slide_number: int, text: str) -> Optional[str]:
        """Regenerate audio for a specific slide with updated text"""
        return await self._generate_slide_audio(text, lesson_id, slide_number)


# Global service instance
_presentation_service: Optional[PresentationService] = None


def get_presentation_service() -> PresentationService:
    """Get or create the presentation service singleton"""
    global _presentation_service
    if _presentation_service is None:
        _presentation_service = PresentationService()
    return _presentation_service
