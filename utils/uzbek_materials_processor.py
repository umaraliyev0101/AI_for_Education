"""
Uzbek Materials Processor
Extracts and processes text from various document formats for the QA system.
"""

import os
import json
from typing import List, Dict, Optional
from pathlib import Path
import logging

# Document processing libraries
import PyPDF2
from pptx import Presentation
from docx import Document

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class UzbekMaterialsProcessor:
    """
    Processes educational materials in various formats (PDF, PPTX, DOCX, TXT)
    and extracts text for the QA system.
    """
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the materials processor.
        
        Args:
            chunk_size: Size of text chunks for splitting
            chunk_overlap: Overlap between consecutive chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_formats = ['.pdf', '.pptx', '.docx', '.txt', '.md']
    
    def extract_from_pdf(self, pdf_path: str) -> str:
        """
        Extract text from a PDF file.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            Extracted text content
        """
        try:
            text = ""
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                logger.info(f"PDF dan matn ajratilmoqda: {pdf_path} ({num_pages} sahifa)")
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Sahifa {page_num} ---\n"
                        text += page_text + "\n"
                        
            logger.info(f"PDF dan {len(text)} belgi ajratildi")
            return text
        except Exception as e:
            logger.error(f"PDF dan matn ajratishda xatolik: {pdf_path} - {str(e)}")
            return ""
    
    def extract_from_pptx(self, pptx_path: str) -> str:
        """
        Extract text from a PowerPoint presentation.
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Extracted text content
        """
        try:
            prs = Presentation(pptx_path)
            text = ""
            logger.info(f"PowerPoint dan matn ajratilmoqda: {pptx_path} ({len(prs.slides)} slayd)")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slayd {slide_num} ---\n"
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            text += row_text + "\n"
                    
                    # Extract text from notes
                    if hasattr(slide, "notes_slide") and slide.notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text:
                            text += f"Izohlar: {notes_text}\n"
            
            logger.info(f"PowerPoint dan {len(text)} belgi ajratildi")
            return text
        except Exception as e:
            logger.error(f"PowerPoint dan matn ajratishda xatolik: {pptx_path} - {str(e)}")
            return ""
    
    def extract_from_docx(self, docx_path: str) -> str:
        """
        Extract text from a Word document.
        
        Args:
            docx_path: Path to the DOCX file
            
        Returns:
            Extracted text content
        """
        try:
            doc = Document(docx_path)
            text = ""
            logger.info(f"Word hujjatidan matn ajratilmoqda: {docx_path}")
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"
            
            logger.info(f"Word hujjatidan {len(text)} belgi ajratildi")
            return text
        except Exception as e:
            logger.error(f"Word hujjatidan matn ajratishda xatolik: {docx_path} - {str(e)}")
            return ""
    
    def extract_from_txt(self, txt_path: str) -> str:
        """
        Extract text from a plain text file.
        
        Args:
            txt_path: Path to the text file
            
        Returns:
            File content
        """
        try:
            with open(txt_path, 'r', encoding='utf-8') as file:
                text = file.read()
            logger.info(f"Matn faylidan {len(text)} belgi o'qildi: {txt_path}")
            return text
        except Exception as e:
            logger.error(f"Matn faylini o'qishda xatolik: {txt_path} - {str(e)}")
            return ""
    
    def extract_text(self, file_path: str) -> str:
        """
        Extract text from a file based on its extension.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text content
        """
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return self.extract_from_pdf(file_path)
        elif ext == '.pptx':
            return self.extract_from_pptx(file_path)
        elif ext == '.docx':
            return self.extract_from_docx(file_path)
        elif ext in ['.txt', '.md']:
            return self.extract_from_txt(file_path)
        else:
            logger.warning(f"Qo'llab-quvvatlanmaydigan fayl formati: {ext}")
            return ""
    
    def split_text_into_chunks(self, text: str) -> List[str]:
        """
        Split text into overlapping chunks for better context preservation.
        
        Args:
            text: Text to split
            
        Returns:
            List of text chunks
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = min(start + self.chunk_size, text_length)
            
            # Try to break at sentence boundary
            if end < text_length:
                # Look for sentence endings
                for sep in ['. ', '? ', '! ', '\n\n', '\n']:
                    last_sep = text.rfind(sep, start, end)
                    if last_sep != -1 and last_sep > start:
                        end = last_sep + len(sep)
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move to next chunk with overlap
            next_start = end - self.chunk_overlap
            if next_start <= start:
                # Prevent infinite loop
                next_start = start + max(1, self.chunk_size // 2)
            start = next_start
        
        logger.info(f"Matn {len(chunks)} qismlarga bo'lingan")
        return chunks
    
    def process_materials(self, file_paths: List[str]) -> List[Dict[str, str]]:
        """
        Process multiple material files and return chunks with metadata.
        
        Args:
            file_paths: List of file paths to process
            
        Returns:
            List of dictionaries containing text chunks and metadata
        """
        all_chunks = []
        
        for file_path in file_paths:
            if not os.path.exists(file_path):
                logger.warning(f"Fayl topilmadi: {file_path}")
                continue
            
            logger.info(f"Fayl qayta ishlanmoqda: {file_path}")
            
            # Extract text
            text = self.extract_text(file_path)
            
            if not text:
                continue
            
            # Split into chunks
            chunks = self.split_text_into_chunks(text)
            
            # Add metadata
            file_name = Path(file_path).name
            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    'text': chunk,
                    'source': file_name,
                    'chunk_id': i,
                    'file_path': file_path
                })
        
        logger.info(f"Jami {len(all_chunks)} qism yaratildi {len(file_paths)} fayldan")
        return all_chunks
    
    def process_directory(self, directory_path: str) -> List[Dict[str, str]]:
        """
        Process all supported files in a directory.
        
        Args:
            directory_path: Path to the directory
            
        Returns:
            List of text chunks with metadata
        """
        file_paths = []
        
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                if Path(file).suffix.lower() in self.supported_formats:
                    file_paths.append(os.path.join(root, file))
        
        logger.info(f"Papkadan {len(file_paths)} fayl topildi: {directory_path}")
        return self.process_materials(file_paths)
    
    def save_processed_materials(self, chunks: List[Dict[str, str]], output_path: str):
        """
        Save processed materials to a JSON file.
        
        Args:
            chunks: Processed text chunks
            output_path: Path to save the JSON file
        """
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(chunks, f, ensure_ascii=False, indent=2)
            logger.info(f"Qayta ishlangan materiallar saqlandi: {output_path}")
        except Exception as e:
            logger.error(f"Faylga saqlashda xatolik: {output_path} - {str(e)}")
    
    def load_processed_materials(self, input_path: str) -> List[Dict[str, str]]:
        """
        Load processed materials from a JSON file.
        
        Args:
            input_path: Path to the JSON file
            
        Returns:
            List of text chunks with metadata
        """
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                chunks = json.load(f)
            logger.info(f"Qayta ishlangan materiallar yuklandi: {input_path} ({len(chunks)} qism)")
            return chunks
        except Exception as e:
            logger.error(f"Fayldan yuklashda xatolik: {input_path} - {str(e)}")
            return []


if __name__ == "__main__":
    # Example usage
    processor = UzbekMaterialsProcessor(chunk_size=500, chunk_overlap=100)
    
    # Process a single file
    # chunks = processor.process_materials(['path/to/your/file.pdf'])
    
    # Process a directory
    # chunks = processor.process_directory('path/to/materials')
    
    # Save processed materials
    # processor.save_processed_materials(chunks, 'processed_materials.json')
    
    print("Uzbek Materials Processor tayyor!")
    print(f"Qo'llab-quvvatlanadigan formatlar: {processor.supported_formats}")
